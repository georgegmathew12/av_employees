"""Generate av_employees_pipeline.pptx — a presentation deck for the ELT pipeline."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- palette ----------
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x2D, 0x6C, 0xDF)
ACCENT_LIGHT = RGBColor(0xE6, 0xEF, 0xFD)
BRONZE = RGBColor(0xB0, 0x7A, 0x3E)
SILVER = RGBColor(0x9C, 0xA3, 0xAF)
GOLD = RGBColor(0xD4, 0xA8, 0x2E)
BG = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def textbox(slide, left, top, width, height, text, *, size=18, bold=False,
            color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, title, kicker=None):
    if kicker:
        textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.35),
                kicker, size=12, bold=True, color=ACCENT)
    textbox(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.7),
            title, size=30, bold=True, color=INK)
    rule = slide.shapes.add_connector(1, Inches(0.6), Inches(1.35), Inches(12.7), Inches(1.35))
    rule.line.color.rgb = RULE
    rule.line.width = Pt(1)


def footer(slide, hr_note=None, page=None):
    if hr_note:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.7),
                                     Inches(12.1), Inches(0.5))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT_LIGHT
        tb = bar.text_frame
        tb.margin_left = Inches(0.2)
        tb.margin_right = Inches(0.2)
        tb.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tb.paragraphs[0]
        r = p.add_run()
        r.text = f"Why this matters for HR: {hr_note}"
        r.font.name = "Calibri"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = ACCENT
    if page is not None:
        textbox(slide, Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3),
                str(page), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, left, top, width, height, items, *, size=18, color=INK,
            bold_first_word=False, gap=Pt(6)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = ACCENT
        r1.font.bold = True
        r1.font.name = "Calibri"
        if bold_first_word and " " in item:
            first, rest = item.split(" ", 1)
            r2 = p.add_run()
            r2.text = first + " "
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.bold = True
            r2.font.name = "Calibri"
            r3 = p.add_run()
            r3.text = rest
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
            r3.font.name = "Calibri"
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
    return tb


def notes(slide, lines):
    nf = slide.notes_slide.notes_text_frame
    nf.text = lines[0] if lines else ""
    for line in lines[1:]:
        p = nf.add_paragraph()
        p.text = line


def pill(slide, left, top, width, height, text, fill, text_color=BG, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.3
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.font.name = "Calibri"
    return box


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(2.25)
    # add arrowhead via XML
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "h": "med"})
    ln.append(tail)
    return conn


def table_simple(slide, left, top, width, height, headers, rows, *,
                 first_col_bold=True, header_fill=ACCENT, header_color=BG):
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table
    # column widths: first col wider
    if ncols == 3:
        tbl.columns[0].width = int(width * 0.28)
        tbl.columns[1].width = int(width * 0.36)
        tbl.columns[2].width = int(width * 0.36)
    elif ncols == 2:
        tbl.columns[0].width = int(width * 0.30)
        tbl.columns[1].width = int(width * 0.70)
    elif ncols == 4:
        for c in range(ncols):
            tbl.columns[c].width = int(width / ncols)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = header_color
        r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if i % 2 == 1 else RGBColor(0xF7, 0xF8, 0xFA)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.size = Pt(12)
            r.font.color.rgb = INK
            r.font.bold = bool(first_col_bold and j == 0)
            r.font.name = "Calibri"
    return tbl


# =========================================================
# Slide 1 — Title
# =========================================================
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(2.4))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
textbox(s, Inches(0.8), Inches(0.7), Inches(12), Inches(0.5),
        "AV EMPLOYEES  /  ELT PIPELINE", size=13, bold=True, color=BG)
textbox(s, Inches(0.8), Inches(1.2), Inches(12), Inches(1.0),
        "From 7 raw CSVs to trusted HR analytics", size=38, bold=True, color=BG)
textbox(s, Inches(0.8), Inches(2.7), Inches(12), Inches(0.8),
        "A governed, tested data pipeline that Tableau and analysts can rely on.",
        size=22, color=INK)
textbox(s, Inches(0.8), Inches(3.6), Inches(12), Inches(0.6),
        "Built with Fivetran  ·  Snowflake  ·  dbt",
        size=18, bold=True, color=MUTED)
textbox(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
        "Prepared by George Mathew  ·  May 2026",
        size=12, color=MUTED)
notes(s, [
    "Open: this deck shows how raw HR CSVs became a governed analytics dataset.",
    "Two audiences — frame results for the HR leader, depth for the analysts.",
    "Total: ~25 min + Q&A.",
])

# =========================================================
# Slide 2 — Business question
# =========================================================
s = add_slide()
header(s, "What HR wanted to answer", kicker="THE BUSINESS QUESTION")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.5), Inches(4.5), [
    "Headcount — how many employees, broken down how?",
    "Attrition — who left, when, and why?",
    "Tenure — how long do people stay?",
    "Departmental composition — who works where?",
    "Exit reasons — what patterns show up?",
], size=20)
# right card
card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.7),
                         Inches(5.1), Inches(4.5))
card.adjustments[0] = 0.05
card.fill.solid(); card.fill.fore_color.rgb = ACCENT_LIGHT
card.line.color.rgb = ACCENT; card.line.width = Pt(1)
textbox(s, Inches(7.8), Inches(1.85), Inches(4.8), Inches(0.5),
        "Example dashboard tiles", size=14, bold=True, color=ACCENT)
bullets(s, Inches(7.8), Inches(2.4), Inches(4.8), Inches(3.8), [
    "Headcount by department",
    "Monthly attrition rate",
    "Tenure distribution at exit",
    "Top exit reasons",
    "Employee count by generation",
], size=16)
notes(s, [
    "Anchor the room in HR outcomes before introducing any tool.",
    "Skip tech jargon entirely on this slide.",
])

# =========================================================
# Slide 3 — Before vs After
# =========================================================
s = add_slide()
header(s, "Before vs. after", kicker="THE TRANSFORMATION")
table_simple(
    s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.2),
    ["", "Before", "After"],
    [
        ["Source of truth", "7 CSVs scattered in Google Drive",
         "One governed star schema in Snowflake"],
        ["Refresh", "Manual file drops",
         "Fivetran-managed ingestion"],
        ["Trust", "\"Which file is current?\"",
         "Tests + contracts on every model"],
        ["Access", "Whoever holds the file",
         "Tableau dashboards + SQL views, role-grantable"],
        ["Audit trail", "None",
         "Full lineage from CSV → dashboard column"],
    ],
)
footer(s, "Every number on a dashboard can be traced back to the raw CSV row.", page=3)
notes(s, [
    "Most important slide for the HR leader — read each row out loud.",
    "Pause on 'Trust' — this is the punchline of the project.",
])

# =========================================================
# Slide 4 — Pipeline diagram
# =========================================================
s = add_slide()
header(s, "The pipeline in one picture", kicker="ARCHITECTURE")

y = Inches(2.6)
h = Inches(1.1)
# six stages
stages = [
    ("CSVs",         "Google Drive",       MUTED,   Inches(0.6),  Inches(1.5)),
    ("Fivetran",     "Managed ingest",     ACCENT,  Inches(2.4),  Inches(1.6)),
    ("Snowflake raw","Landing schema",     INK,     Inches(4.3),  Inches(1.8)),
    ("dbt",          "bronze → silver → gold", ACCENT, Inches(6.4), Inches(2.2)),
    ("Snowflake gold","Star schema + view", INK,    Inches(9.0),  Inches(1.9)),
    ("Tableau / SQL","Consumers",          GOLD,    Inches(11.1), Inches(1.8)),
]
positions = []
for label, sub, color, left, width in stages:
    pill(s, left, y, width, h, label, color, size=15)
    textbox(s, left, y + h + Inches(0.1), width, Inches(0.4),
            sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    positions.append((left, width))

# arrows between stages
for i in range(len(positions) - 1):
    l1, w1 = positions[i]
    l2, _ = positions[i + 1]
    arrow(s, l1 + w1, y + h / 2, l2, y + h / 2)

# bottom legend
textbox(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
        "Each tool has one job. The seams between them are the testable contracts.",
        size=16, color=INK, bold=True)
bullets(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.2), [
    "Fivetran — moves bytes. Not a transformer.",
    "Snowflake — stores everything; compute scales separately from storage.",
    "dbt — every transformation lives in version-controlled, tested SQL.",
], size=14)
footer(s, page=4)
notes(s, [
    "Reuse this diagram on the medallion and quality slides.",
    "Stress: each tool has one job. No bespoke loaders or scripts to maintain.",
])

# =========================================================
# Slide 5 — Why these tools
# =========================================================
s = add_slide()
header(s, "Why these three tools", kicker="DECISIONS")
table_simple(
    s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.6),
    ["Tool", "Why chosen", "Trade-off"],
    [
        ["Fivetran", "Managed CSV → Snowflake ingestion. Zero loader code to maintain.",
         "Per-row cost. Justified for a recurring source, not a one-off."],
        ["Snowflake", "Storage and compute scale separately. SQL-native — analysts self-serve.",
         "Cloud cost vs. on-prem. Wins on time-to-value and reliability."],
        ["dbt", "Version-controlled SQL transforms, built-in tests, contracts, lineage.",
         "Replaces \"someone's spreadsheet.\" Learning curve for non-engineers."],
    ],
)
footer(s, "Off-the-shelf tools mean we spend time on the data, not on plumbing.",
       page=5)
notes(s, [
    "If asked about cost: small footprint here — Fivetran rows are bounded, Snowflake bills only when dbt runs.",
])

# =========================================================
# Slide 6 — Medallion architecture
# =========================================================
s = add_slide()
header(s, "Bronze → Silver → Gold", kicker="HOW DBT IS ORGANIZED")

# three cards
def layer_card(left, color, title, subtitle, items):
    w = Inches(4.0)
    y = Inches(1.8)
    h = Inches(4.4)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid(); card.fill.fore_color.rgb = BG
    card.line.color.rgb = color; card.line.width = Pt(2)
    # header band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, Inches(0.8))
    band.adjustments[0] = 0.10
    band.fill.solid(); band.fill.fore_color.rgb = color
    band.line.fill.background()
    textbox(s, left, y + Inches(0.15), w, Inches(0.5),
            title, size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)
    textbox(s, left + Inches(0.2), y + Inches(0.95), w - Inches(0.4), Inches(0.5),
            subtitle, size=12, bold=True, color=color)
    bullets(s, left + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(2.8),
            items, size=13)

layer_card(Inches(0.6),  BRONZE, "BRONZE", "Typed raw — 1:1 with source", [
    "Renames columns to project standard",
    "Keeps loaded_at + _line tiebreaker",
    "No filtering, no business logic",
    "One model per source table",
])
layer_card(Inches(4.85), SILVER, "SILVER", "Cleaned + conformed", [
    "Drops null rows, parses VARCHAR dates",
    "Dedupes by loaded_at then _line",
    "Removes FK orphans",
    "Joins related entities (employee, dept)",
])
layer_card(Inches(9.1),  GOLD,   "GOLD",   "Star schema + flat view", [
    "Dim + fact tables with contracts",
    "fct_employment, fct_employee_department",
    "vw_employee_full for Excel / Snowsight",
    "What Tableau and analysts query",
])
footer(s, "Three checkpoints. If a number looks wrong, we trace it back through audited stops.",
       page=6)
notes(s, [
    "Use the metaphor: bronze = raw ingredients, silver = prepped, gold = plated.",
])

# =========================================================
# Slide 7 — Gold layer detail
# =========================================================
s = add_slide()
header(s, "The gold layer — what HR actually uses", kicker="DATA MART")

# Left: tables
textbox(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.5),
        "Star schema", size=18, bold=True, color=ACCENT)
bullets(s, Inches(0.7), Inches(2.2), Inches(6.0), Inches(3.5), [
    "fct_employment — 1 row per employee",
    "fct_employee_department — bridge (multi-dept)",
    "dim_employee — name, gender, hire date, generation",
    "dim_department — department names",
    "dim_title — job titles",
    "dim_exit_reason — code + label",
    "dim_date — calendar dim",
    "dim_generation — Boomer / Gen X / Millennial / Gen Z",
], size=14)

# Right: flat view
textbox(s, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.5),
        "Flat view for non-Tableau users", size=18, bold=True, color=ACCENT)
card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.2),
                         Inches(5.5), Inches(3.5))
card.adjustments[0] = 0.05
card.fill.solid(); card.fill.fore_color.rgb = ACCENT_LIGHT
card.line.color.rgb = ACCENT; card.line.width = Pt(1)
textbox(s, Inches(7.4), Inches(2.4), Inches(5.2), Inches(0.5),
        "vw_employee_full", size=16, bold=True, color=INK, font="Consolas")
textbox(s, Inches(7.4), Inches(2.9), Inches(5.2), Inches(2.7),
        "One row per employee, all attributes pre-joined. "
        "Departments collapsed to a comma-separated string. "
        "Open it in Excel, Snowsight, or any SQL client — no joins required.",
        size=13, color=INK)
footer(s, "Analysts get the star schema. Everyone else gets one big table.",
       page=7)
notes(s, [
    "This is the slide HR remembers — show a screenshot of vw_employee_full if you can.",
])

# =========================================================
# Slide 8 — Quality controls
# =========================================================
s = add_slide()
header(s, "How we know the data is right", kicker="QUALITY & TRUST")
bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.5), [
    "Tests on every layer — unique, not_null, accepted_values, relationships",
    "Contracts on gold — a schema change breaks the build, not the dashboard",
    "Snapshots — capture row-level history for tables with no source dates",
    "CI — every code change is parsed before it can merge",
    "Lineage docs — click any gold column and see exactly where it came from",
], size=20, bold_first_word=True)
footer(s, "If the data breaks, the pipeline fails loudly — it does not silently publish bad numbers.",
       page=8)
notes(s, [
    "The trust slide. Slow down. Each bullet = one anxiety the HR leader has.",
])

# =========================================================
# Slide 9 — What we gained
# =========================================================
s = add_slide()
header(s, "What was gained", kicker="OUTCOMES")
bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0), [
    "Single source of truth for headcount and attrition",
    "Reproducible — anyone can rebuild from raw in minutes",
    "Self-serve for analysts (star schema) and non-technical users (flat view)",
    "Full lineage — click a column, see where it came from",
    "Governed access — Snowflake roles control who sees what",
    "Tested — broken data is caught before it reaches a dashboard",
], size=20)
footer(s, page=9)
notes(s, [
    "Frame as time-savings and trust, not technology.",
])

# =========================================================
# Slide 10 — Limitations
# =========================================================
s = add_slide()
header(s, "What was lost — honest limitations", kicker="DATA GAPS (ALL UPSTREAM)")
table_simple(
    s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(5.0),
    ["Gap", "Effect on HR", "What unlocks it"],
    [
        ["No dates on dept assignments", "Can't say which dept at exit; can't measure manager span",
         "Source adds dates to dept_emp / dept_manager"],
        ["exit_reason is a code only", "Showing raw codes, not labels",
         "Source supplies a decoder table"],
        ["No location column anywhere", "Can't slice headcount or attrition by location",
         "Source adds location data"],
        ["Salary & title — 1 row each, no dates", "Treated as \"current\" — no comp or promotion trends",
         "Source adds dated history records"],
        ["~77% of salary / dept rows have unknown IDs", "Those rows are dropped from analysis",
         "Source fixes referential integrity"],
        ["Dates stored as 2-digit-year text", "Ambiguity around year boundaries",
         "Source supplies ISO dates"],
    ],
)
footer(s, "Every gap is upstream — fixing the source unlocks new analysis without rebuilding the pipeline.",
       page=10)
notes(s, [
    "Be unapologetic about limitations — honesty builds trust.",
    "Each row is a future win waiting on the source team.",
])

# =========================================================
# Slide 11 — Key decisions
# =========================================================
s = add_slide()
header(s, "Key decisions and their impact", kicker="DESIGN CHOICES")
table_simple(
    s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(4.8),
    ["Decision", "Why", "Impact"],
    [
        ["Truncate-and-reload (not incremental)",
         "One-time load; simpler operationally",
         "Fast to ship. Revisit when pipeline becomes recurring."],
        ["Drop FK orphans via inner join in silver",
         "Cleaner gold numbers; avoid null-stuffed joins",
         "~77% of salary rows excluded. Documented gap."],
        ["Bridge fact for departments",
         "Source has multi-dept employees with no dates",
         "Correct multi-dept counts; flag for single-dept lookups."],
        ["Flat view AND star schema",
         "Different consumers — Tableau vs. Excel / Snowsight",
         "Unblocks non-technical users at low maintenance cost."],
        ["Snapshots on undated tables",
         "Capture history from this day forward",
         "Future-proofs analysis even before source adds dates."],
    ],
)
footer(s, page=11)
notes(s, [
    "Analysts: this is where they'll have questions. Be ready to defend the inner-join call.",
])

# =========================================================
# Slide 12 — Roadmap
# =========================================================
s = add_slide()
header(s, "Roadmap — ordered by HR ROI", kicker="WHAT'S NEXT")
bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0), [
    "Schedule recurring sync — Fivetran + dbt on a daily cadence",
    "Salary history fact — enables comp trends, raises, growth analysis",
    "Title history fact — promotion paths, time-in-role",
    "Dated department assignments — manager span, hiring velocity, dept-at-exit",
    "Location dimension — if source ever supplies it",
    "Source freshness alerts — catch upstream failures fast",
    "Tableau role grants in Snowflake — analyst_role provisioning",
], size=18, bold_first_word=True)
footer(s, "The biggest wins come from source-data improvements — pipeline is ready to absorb them.",
       page=12)
notes(s, [
    "Order by HR value, not engineering effort.",
])

# =========================================================
# Slide 13 — Future charts
# =========================================================
s = add_slide()
header(s, "Charts we could add once source improves", kicker="FUTURE DASHBOARDS")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(5.0), [
    "Salary growth over time",
    "Promotion paths + relation to exits",
    "Average time in role before promotion",
    "Manager span of control",
], size=17)
bullets(s, Inches(7.0), Inches(1.7), Inches(6.0), Inches(5.0), [
    "Hiring velocity by department",
    "Cohort retention by hire year",
    "Compensation vs. tenure curves",
    "Department-at-exit trends",
], size=17)
footer(s, "These are the questions HR will ask next quarter — pipeline is ready.",
       page=13)
notes(s, [
    "Frames the next HR conversation. Each bullet has a clear source-data unlock.",
])

# =========================================================
# Slide 14 — Anticipated questions
# =========================================================
s = add_slide()
header(s, "Anticipated questions", kicker="Q&A PREP")
textbox(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(0.5),
        "From HR leadership", size=15, bold=True, color=ACCENT)
bullets(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(4.6), [
    "Can I trust these numbers?",
    "How fresh is the data?",
    "What if an upstream column changes?",
    "Who can see this data?",
    "What does this cost to run?",
    "What about GDPR / employee deletes?",
    "How do we add a new system (payroll, HRIS)?",
], size=13)
textbox(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.5),
        "From analysts", size=15, bold=True, color=ACCENT)
bullets(s, Inches(7.0), Inches(2.1), Inches(6.0), Inches(4.6), [
    "What is the grain of each fact?",
    "How are slowly-changing dimensions handled?",
    "Why inner-join orphans instead of left-join + null?",
    "Test coverage — what tiers of CI?",
    "How are 2-digit years disambiguated?",
    "Why a flat view AND a star schema?",
    "Snapshot strategy for undated source tables?",
], size=13)
footer(s, page=14)
notes(s, [
    "Don't read these out — keep as backup. Rehearse the answers in advance.",
])

# =========================================================
# Slide 15 — Close
# =========================================================
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = ACCENT
textbox(s, Inches(0.8), Inches(0.8), Inches(12), Inches(0.5),
        "CLOSING", size=13, bold=True, color=BG)
textbox(s, Inches(0.8), Inches(1.3), Inches(12), Inches(1.5),
        "What can HR answer now that it couldn't before?",
        size=34, bold=True, color=BG)
textbox(s, Inches(0.8), Inches(3.0), Inches(12), Inches(2.5),
        "Every headcount, attrition, and tenure number is now reproducible, "
        "tested, and traceable to a row of raw data.\n\n"
        "The pipeline absorbs source-data improvements without rebuilding — "
        "every gap we identified today becomes a new chart tomorrow.",
        size=20, color=BG)
textbox(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.5),
        "Thank you — questions?",
        size=18, bold=True, color=BG)
notes(s, [
    "End on the dashboard demo if time allows.",
    "Hand out the gaps + roadmap 1-pager.",
])


out = "/Users/georgemathew/Downloads/personal/av_employees/av_employees_pipeline.pptx"
prs.save(out)
print(f"wrote {out}")
